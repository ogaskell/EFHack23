"""Contains a number of classes used to parse a given file/s, to retrieve slides and notes."""

from abc import ABC, abstractmethod
from typing import Any, Optional

import pptx


class Presentation:
    """Prepresents the combined slides and notes."""

    def __init__(self):
        self.slides: Any = None  # Not currently used
        self.notes: list[list[str]] = []  # notes[slide][bullet_point]


class BaseParser(ABC):
    """Base class for document parsers."""

    presentation: Optional[Presentation] = None

    @abstractmethod
    def load(*args, **kwargs) -> None:
        """Load some file/s to generate a Presentation object."""

    def get_presentation(self) -> Presentation:
        """Will return a Presentation object."""
        if self.presentation is None:
            raise ValueError("Presentation not loaded. Call .load() on the parser first.")

        return self.presentation


class MarkdownParser(BaseParser):
    """Doesn't load any slides, just loads a .md file of notes.

    Uses headings to delimit slides, and bulleted lists as notes.
    """

    def load(self, filename: str) -> None:
        """Load a markdown file."""
        self.presentation = Presentation()
        self.presentation.notes = []

        f = open(filename, "r")
        for raw_line in f:
            line = raw_line.strip()
            if len(line) == 0:
                continue
            elif line[0] == "#":
                self.presentation.notes.append([])
            else:
                self.presentation.notes[-1].append(line.lstrip("- \t"))


class PPTXParser(BaseParser):
    """Parses a PPTX File."""

    def load(self, filename: str) -> None:
        """Load a PPTX file."""
        prs = pptx.Presentation(filename)

        self.presentation = Presentation()
        self.presentation.notes = [
            self.process_notes(slide.notes_slide.notes_placeholder.text)
            if slide.has_notes_slide else "" for slide in prs.slides
        ]
        self.presentation.slides = prs.slides

    def process_notes(self, notes: str) -> list[str]:
        """Split notes text into individual points."""
        return [line.strip().lstrip("- \t") for line in notes.strip().split("\n") if line.strip()]
